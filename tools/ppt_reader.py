from pptx import Presentation

def extract_ppt_text(ppt_path):
    prs = Presentation(ppt_path)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            # 1. Extract standard text boxes
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        text += paragraph.text + "\n"
            
            # 2. Extract tables from the slide
            elif getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame.text.strip():
                            text += cell.text_frame.text.strip() + " | "
                    text += "\n"

    return text
